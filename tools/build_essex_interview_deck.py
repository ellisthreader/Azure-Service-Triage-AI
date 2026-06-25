from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"
SCREENSHOTS = OUT_DIR / "screenshots"
PPTX_OUT = OUT_DIR / "essex_ai_ml_engineer_interview.pptx"
NOTES_OUT = OUT_DIR / "essex_ai_ml_engineer_speaker_notes.md"


ESSEX_RED = RGBColor(230, 0, 45)
INK = RGBColor(18, 24, 38)
SLATE = RGBColor(75, 85, 99)
MUTED = RGBColor(107, 114, 128)
LINE = RGBColor(226, 232, 240)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(21, 44, 78)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(202, 138, 4)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color=WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def text(slide, x, y, w, h, value, size=18, bold=False, color=INK, align=None, all_caps=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = value.upper() if all_caps else value
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return box


def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def pill(slide, x, y, w, h, label, fill, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.35
    text(slide, x + 0.08, y + 0.07, w - 0.16, h - 0.11, label, size=9.2, bold=True, color=color, align=PP_ALIGN.CENTER)
    return shape


def header(slide, idx, title, subtitle=None, dark=False):
    color = WHITE if dark else INK
    sub_color = rgb(203, 213, 225) if dark else MUTED
    text(slide, 0.55, 0.32, 0.55, 0.24, f"{idx:02}", size=9, bold=True, color=ESSEX_RED)
    text(slide, 1.0, 0.24, 8.4, 0.48, title, size=24, bold=True, color=color)
    if subtitle:
        text(slide, 1.02, 0.78, 8.2, 0.28, subtitle, size=10.8, color=sub_color)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.17), Inches(12.22), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(226, 232, 240) if not dark else rgb(71, 85, 105)
    line.line.fill.background()


def footer(slide, idx, dark=False):
    color = rgb(203, 213, 225) if dark else MUTED
    text(slide, 0.55, 7.13, 7.6, 0.22, "Service Priority AI | Azure MLOps interview presentation", size=8.2, color=color)
    text(slide, 12.2, 7.13, 0.55, 0.22, str(idx), size=8.2, color=color, align=PP_ALIGN.RIGHT)


def image_cover(slide, path, x, y, w, h, radius=False):
    path = Path(path)
    if not path.exists():
        return None
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    image_ratio = iw / ih
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif image_ratio < box_ratio:
        crop = (1 - image_ratio / box_ratio) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    if radius:
        border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        border.fill.background()
        border.line.color.rgb = LINE
        border.line.width = Pt(1)
        border.adjustments[0] = 0.04
    return pic


def card(slide, x, y, w, h, title, body=None, accent=ESSEX_RED, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    shape.adjustments[0] = 0.05
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, x + 0.22, y + 0.18, w - 0.44, 0.26, title, size=13.2, bold=True, color=INK)
    if body:
        bullets(slide, x + 0.22, y + 0.56, w - 0.45, h - 0.68, body, size=11.5, color=SLATE, gap=3)
    return shape


def metric(slide, x, y, w, label, value, note, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.adjustments[0] = 0.05
    text(slide, x + 0.16, y + 0.12, w - 0.32, 0.2, label, size=8.8, bold=True, color=MUTED, all_caps=True)
    text(slide, x + 0.16, y + 0.34, w - 0.32, 0.34, value, size=21, bold=True, color=accent)
    text(slide, x + 0.16, y + 0.73, w - 0.32, 0.18, note, size=8.5, color=MUTED)


def flow_node(slide, x, y, w, h, title, note, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(248, 250, 252)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.06
    text(slide, x + 0.12, y + 0.13, w - 0.24, 0.25, title, size=11.5, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(slide, x + 0.14, y + 0.48, w - 0.28, 0.38, note, size=8.8, color=SLATE, align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(148, 163, 184)
    conn.line.width = Pt(1.2)


def build_deck():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = ESSEX_RED
    top.line.fill.background()
    text(slide, 0.8, 0.82, 5.7, 0.3, "AI / ML ENGINEER INTERVIEW", size=11, bold=True, color=ESSEX_RED)
    text(slide, 0.78, 2.18, 9.35, 0.72, "Service Priority AI", size=38, bold=True, color=INK)
    text(slide, 0.82, 3.02, 9.55, 0.46, "Responsible Azure MLOps for public-service case prioritisation", size=18, color=SLATE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(4.35), Inches(3.05), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ESSEX_RED
    line.line.fill.background()
    text(slide, 0.82, 5.28, 4.9, 0.28, "Ellis", size=13, bold=True, color=INK)
    text(slide, 0.82, 5.68, 5.4, 0.28, "30 June 2026 | Virtual interview", size=11.5, color=MUTED)
    text(slide, 11.8, 6.95, 0.55, 0.22, "1", size=8.2, color=MUTED, align=PP_ALIGN.RIGHT)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    header(slide, 2, "Context And My Role", "Why the system exists, who uses it, and what I personally built")
    image_cover(slide, SCREENSHOTS / "current-home.png", 7.15, 1.55, 5.0, 3.75, radius=True)
    card(slide, 0.72, 1.55, 3.0, 4.55, "Business problem", [
        "Urgent cases can be hidden in mixed service queues.",
        "Officers need consistent triage support, not an automated decision maker.",
        "Managers need visibility of risk, demand, overrides, and model behaviour.",
    ], accent=ESSEX_RED, fill=PAPER)
    card(slide, 3.95, 1.55, 2.9, 4.55, "Users", [
        "Case officers reviewing service requests.",
        "Service managers monitoring queues and exceptions.",
        "Information Governance and responsible AI reviewers.",
    ], accent=BLUE, fill=PAPER)
    card(slide, 7.15, 5.58, 5.0, 0.85, "My contribution", [
        "Designed and built the ML lifecycle, API boundary, Azure deployment assets, monitoring outputs, and governance documentation.",
    ], accent=GREEN, fill=WHITE)
    footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide, rgb(250, 250, 250))
    header(slide, 3, "Data And Real-World Limitations", "The data was useful for proving the pattern, but the risks are the point")
    metric(slide, 0.72, 1.45, 2.25, "Training rows", "20k", "synthetic service cases", NAVY)
    metric(slide, 3.2, 1.45, 2.25, "Validation rows", "5k", "hold-out evaluation", BLUE)
    metric(slide, 5.68, 1.45, 2.25, "Decision role", "Advisory", "human review required", ESSEX_RED)
    card(slide, 0.72, 2.95, 3.65, 2.92, "Inputs", [
        "Service type, SLA, district, channel, days open, previous contacts.",
        "Vulnerability flag, deprivation band, accessibility need.",
        "Urgency notes transformed with explainable text features.",
    ], accent=BLUE, fill=WHITE)
    card(slide, 4.62, 2.95, 3.65, 2.92, "Limitations", [
        "Synthetic labels do not prove live resident-service performance.",
        "Real notes would be sensitive, noisy, incomplete, and unevenly written.",
        "Vulnerability and deprivation can create proxy-risk concerns.",
    ], accent=AMBER, fill=WHITE)
    card(slide, 8.52, 2.95, 3.65, 2.92, "Controls", [
        "DPIA-lite, model card, responsible AI assessment.",
        "No raw real resident data in the demo.",
        "Fairness slices, drift review, override monitoring, and human accountability.",
    ], accent=GREEN, fill=WHITE)
    footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    header(slide, 4, "Model And Evaluation", "Transparent baseline first, because missed urgent cases are the service risk")
    image_cover(slide, SCREENSHOTS / "live-case-details.png", 7.0, 1.45, 5.15, 3.85, radius=True)
    text(slide, 0.72, 1.52, 5.35, 0.4, "Baseline model", size=20, bold=True, color=INK)
    bullets(slide, 0.74, 2.0, 5.4, 1.45, [
        "scikit-learn LogisticRegression pipeline.",
        "ColumnTransformer over categorical, numeric, boolean, and TF-IDF text features.",
        "Chosen for repeatability, explainability, and public-sector accountability.",
    ], size=14.2, color=SLATE)
    metric(slide, 0.72, 3.85, 1.7, "Accuracy", "90.2%", "validation", BLUE)
    metric(slide, 2.62, 3.85, 1.7, "Macro F1", "90.35%", "balanced view", GREEN)
    metric(slide, 4.52, 3.85, 1.7, "High recall", "94.08%", "harm-oriented", ESSEX_RED)
    card(slide, 0.72, 5.08, 5.5, 1.42, "Validation checks", [
        "Hold-out metrics, confusion matrix, class recall, confidence review, fairness slices.",
        "High-priority recall gate passed; fairness remains review_required.",
    ], accent=NAVY, fill=PAPER)
    footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide, rgb(15, 23, 42))
    header(slide, 5, "Azure Architecture And Service Choices", "Governed serving path with online and batch scoring", dark=True)
    y = 1.55
    nodes = [
        (0.75, y, 1.45, "Synthetic data", "approved demo data", ESSEX_RED),
        (2.55, y, 1.45, "Azure ML pipeline", "train / evaluate / gate", BLUE),
        (4.35, y, 1.45, "Model registry", "versioned artifact", GREEN),
        (6.15, y, 1.6, "Online endpoint", "real-time scoring", GREEN),
        (8.15, y, 1.75, "FastAPI wrapper", "validation + audit", AMBER),
        (10.25, y, 1.45, "React dashboard", "officer workflow", ESSEX_RED),
    ]
    for x, yy, ww, title_, note_, color in nodes:
        flow_node(slide, x, yy, ww, 0.94, title_, note_, color)
    for x1, x2 in [(2.2, 2.55), (4.0, 4.35), (5.8, 6.15), (7.75, 8.15), (9.9, 10.25)]:
        arrow(slide, x1, y + 0.47, x2, y + 0.47)
    flow_node(slide, 6.15, 3.13, 1.6, 0.94, "Batch endpoint", "offline review queues", GREEN)
    flow_node(slide, 8.15, 3.13, 1.75, 0.94, "Audit store", "decisions + overrides", AMBER)
    flow_node(slide, 10.25, 3.13, 1.45, 0.94, "Power BI", "monitoring review", BLUE)
    arrow(slide, 6.95, 2.5, 6.95, 3.13)
    arrow(slide, 7.75, 3.6, 8.15, 3.6)
    arrow(slide, 9.9, 3.6, 10.25, 3.6)
    card(slide, 0.75, 4.95, 3.45, 1.2, "Why Azure ML", [
        "Pipelines, pinned environments, registry, managed endpoints, traffic control, and reproducibility.",
    ], accent=BLUE, fill=rgb(248, 250, 252))
    card(slide, 4.55, 4.95, 3.45, 1.2, "Why API wrapper", [
        "Browser never handles Azure ML endpoint secrets; the API centralises validation, logging, and audit.",
    ], accent=AMBER, fill=rgb(248, 250, 252))
    card(slide, 8.35, 4.95, 3.45, 1.2, "Why batch scoring", [
        "Supports scheduled review queues, monitoring backfills, and retrospective evaluation.",
    ], accent=GREEN, fill=rgb(248, 250, 252))
    footer(slide, 5, dark=True)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    header(slide, 6, "Production Safety Controls", "The system is designed to be supportable before it is impressive")
    image_cover(slide, SCREENSHOTS / "live-signed-in.png", 0.72, 1.45, 5.85, 3.9, radius=True)
    card(slide, 6.9, 1.45, 2.55, 3.9, "Engineering", [
        "Typed FastAPI and Pydantic contracts.",
        "Backend tests, frontend build check, local validation script.",
        "Versioned artifacts, registry tags, and model metadata.",
    ], accent=BLUE, fill=PAPER)
    card(slide, 9.75, 1.45, 2.55, 3.9, "Governance", [
        "Model card, DPIA-lite, responsible AI assessment.",
        "Human review, override reasons, and audit trail.",
        "Fairness status remains review_required before real deployment.",
    ], accent=ESSEX_RED, fill=PAPER)
    card(slide, 0.72, 5.66, 11.58, 0.74, "Operational controls", [
        "Azure Table Storage audit pattern, Application Insights/Azure Monitor health signals, endpoint rollback by traffic/deployment, small compute, scale-to-zero, and monthly budget alerts.",
    ], accent=GREEN, fill=WHITE)
    footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide, rgb(249, 250, 251))
    header(slide, 7, "Monitoring, Incident And Learning", "How I would know the service was degrading and what changed after a rollout issue")
    card(slide, 0.72, 1.45, 3.55, 4.65, "Monitoring", [
        "API health, request volume, error rate, p95 latency, model-loaded state.",
        "Class mix, confidence, low-confidence rate, explanation coverage.",
        "Drift, fairness slices, high-priority recall, audit completeness, and overrides.",
    ], accent=BLUE, fill=WHITE)
    card(slide, 4.62, 1.45, 3.55, 4.65, "Incident", [
        "Linux App Service plan creation in UK South was blocked by subscription VM quota.",
        "Root cause was regional quota, not model code or schema.",
        "No residents were affected; this was synthetic portfolio telemetry.",
    ], accent=ESSEX_RED, fill=WHITE)
    card(slide, 8.52, 1.45, 3.55, 4.65, "Change afterwards", [
        "Used Azure Functions Flex Consumption for the browser API wrapper.",
        "Kept Azure ML behind a governed layer.",
        "Documented fallback hosting path and evidence checks before claiming readiness.",
    ], accent=GREEN, fill=WHITE)
    footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide, rgb(15, 23, 42))
    header(slide, 8, "What I Would Bring To ECC", "Azure-native delivery with responsible AI controls by design", dark=True)
    bullets(slide, 0.82, 1.55, 5.6, 3.0, [
        "Move experiments into governed Azure ML pipelines and model promotion gates.",
        "Deploy through managed online endpoints, batch endpoints, and supportable API boundaries.",
        "Work with Information Governance before real data: DPIA, UK GDPR, access control, retention, and audit.",
        "Coach analysts and data scientists on reproducible, testable MLOps workflows.",
        "Surface monitoring through Power BI-ready outputs that service owners can actually use.",
    ], size=16, color=rgb(226, 232, 240), gap=6)
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.55), Inches(4.75), Inches(2.65))
    quote.fill.solid()
    quote.fill.fore_color.rgb = rgb(30, 41, 59)
    quote.line.color.rgb = rgb(71, 85, 105)
    quote.adjustments[0] = 0.04
    text(slide, 7.42, 1.98, 4.0, 1.55, "Deployment is not just getting a model behind an endpoint. It is making it reproducible, explainable, monitored, governed, and safe to challenge.", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(slide, 7.05, 4.58, 4.75, 1.42, "LLM use disclosure", [
        "I used an AI assistant to structure and polish the deck.",
        "I verified the content against repository artifacts, model metrics, Azure notes, and screenshots from the running website.",
    ], accent=ESSEX_RED, fill=rgb(248, 250, 252))
    footer(slide, 8, dark=True)

    prs.save(PPTX_OUT)


def build_notes():
    notes = """# Essex AI / ML Engineer Presentation Speaker Notes

Target: 8 slides in 10 minutes. The updated deck uses screenshots captured from the currently running local website on 2026-06-22.

## 1. Responsible Azure MLOps
Open by setting the boundary clearly: this is a real cloud-oriented deployment pattern using synthetic public-sector data. The model is advisory and supports human review. The point is safe delivery: reproducibility, governed serving, monitoring, audit, and responsible AI documentation.

## 2. Context And My Role
Explain the operational problem: urgent service requests can sit inside mixed queues, and managers need evidence of how decisions are being supported. State your contribution directly: model lifecycle, API boundary, Azure deployment assets, monitoring outputs, and governance documentation.

## 3. Data And Real-World Limitations
Mention 20,000 training rows and 5,000 validation rows. Then focus on the data risks: synthetic data does not prove live performance, vulnerability fields can be incomplete, deprivation bands can be proxy-risk fields, and free text can encode access and language bias. Your mitigation is governance and human review, not pretending the data is perfect.

## 4. Model And Evaluation
Describe the transparent baseline: a scikit-learn LogisticRegression pipeline over structured and TF-IDF text features. Call out the metrics: 90.2% accuracy, 90.35% macro F1, and 94.08% high-priority recall. Explain that high-priority recall matters because missed urgent cases are a higher service risk.

## 5. Azure Architecture And Service Choices
Walk left to right. Azure ML handles training, environments, registry, online endpoint, and batch endpoint. The browser calls a FastAPI wrapper rather than Azure ML directly. The wrapper centralises validation, logging, audit, and auth boundaries. Batch scoring supports review queues and monitoring backfills.

## 6. Production Safety Controls
Tie this slide to the job description. Engineering controls provide repeatability. Governance controls provide explainability, auditability, and UK GDPR-aware delivery. Operational controls make the system supportable: audit store, Azure Monitor/Application Insights, rollback, small compute, scale-to-zero, and budget alerts.

## 7. Monitoring, Incident And Learning
Start with what you monitor: health, latency, errors, class mix, confidence, drift, fairness, explanation coverage, audit completeness, and overrides. Then describe the rollout issue: UK South Linux App Service quota blocked App Service plan creation. Mitigation was Azure Functions Flex Consumption for the browser API wrapper while keeping Azure ML behind a governed layer.

## 8. What I Would Bring To ECC
Close by mapping to ECC: Azure ML pipelines, managed online endpoints, batch endpoints, responsible AI, Information Governance, Power BI reporting, and coaching colleagues. Include the LLM disclosure exactly as requested in the interview email.
"""
    NOTES_OUT.write_text(notes, encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    build_notes()
    print(PPTX_OUT)
    print(NOTES_OUT)
